#!/usr/bin/env python3
"""
NEVPAZ Scientific Visualization Toolkit
========================================
Hochwertige wissenschaftliche Diagramme, Grafiken und Visualisierungen
für medizinische und neurowissenschaftliche Daten.

Verwendung:
    python scientific_viz.py --demo          # Erstellt Demo-Grafiken
    python scientific_viz.py --type forest   # Erstellt einen Forest Plot
    python scientific_viz.py --type pathway  # Erstellt ein Pathway-Diagramm

Abhängigkeiten:
    pip install matplotlib seaborn plotly numpy pandas python-pptx kaleido Pillow
"""

import argparse
import os
import sys
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle
import matplotlib.patheffects as pe
from matplotlib.collections import PatchCollection
import numpy as np

try:
    import seaborn as sns
    HAS_SEABORN = True
except ImportError:
    HAS_SEABORN = False

try:
    import plotly.graph_objects as go
    import plotly.express as px
    from plotly.subplots import make_subplots
    HAS_PLOTLY = True
except ImportError:
    HAS_PLOTLY = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# =============================================================================
# Farbpaletten
# =============================================================================

COLORS = {
    # NEVPAZ Branding
    'primary_dark': '#1a5276',
    'primary': '#2980b9',
    'primary_light': '#5dade2',
    'accent': '#F39C12',

    # Neurotransmitter
    'dopamin': '#E74C3C',
    'serotonin': '#3498DB',
    'noradrenalin': '#E67E22',
    'gaba': '#2ECC71',
    'glutamat': '#9B59B6',
    'acetylcholin': '#1ABC9C',

    # Hirnregionen
    'prefrontal_cortex': '#3498DB',
    'basal_ganglia': '#E74C3C',
    'hippocampus': '#2ECC71',
    'amygdala': '#E67E22',
    'thalamus': '#9B59B6',
    'hypothalamus': '#F39C12',
    'brainstem': '#1ABC9C',
    'cerebellum': '#34495E',
    'nucleus_accumbens': '#E91E63',
    'vta': '#FF5722',
    'substantia_nigra': '#795548',
    'locus_coeruleus': '#FF9800',
    'raphe_nuclei': '#03A9F4',

    # Allgemein
    'significant': '#27AE60',
    'not_significant': '#E74C3C',
    'background': '#FFFFFF',
    'background_alt': '#F8F9FA',
    'grid': '#E0E0E0',
    'text_dark': '#2C3E50',
    'text_light': '#7F8C8D',
}

NEUROTRANSMITTER_COLORS = {
    'Dopamin': COLORS['dopamin'],
    'Serotonin': COLORS['serotonin'],
    'Noradrenalin': COLORS['noradrenalin'],
    'GABA': COLORS['gaba'],
    'Glutamat': COLORS['glutamat'],
    'Acetylcholin': COLORS['acetylcholin'],
}


def _setup_style():
    """Konfiguriert den globalen Matplotlib-Stil für wissenschaftliche Grafiken."""
    plt.rcParams.update({
        'font.family': 'sans-serif',
        'font.sans-serif': ['Calibri', 'Arial', 'DejaVu Sans', 'Helvetica'],
        'font.size': 12,
        'axes.titlesize': 18,
        'axes.titleweight': 'bold',
        'axes.labelsize': 13,
        'axes.labelcolor': COLORS['text_dark'],
        'axes.edgecolor': COLORS['grid'],
        'axes.facecolor': COLORS['background'],
        'axes.grid': True,
        'axes.spines.top': False,
        'axes.spines.right': False,
        'grid.color': COLORS['grid'],
        'grid.linestyle': '--',
        'grid.alpha': 0.5,
        'figure.facecolor': COLORS['background'],
        'figure.dpi': 150,
        'savefig.dpi': 300,
        'savefig.bbox': 'tight',
        'savefig.facecolor': COLORS['background'],
        'legend.framealpha': 0.9,
        'legend.edgecolor': COLORS['grid'],
        'xtick.color': COLORS['text_dark'],
        'ytick.color': COLORS['text_dark'],
    })


_setup_style()


# =============================================================================
# Forest Plot
# =============================================================================

def create_forest_plot(
    studies: list,
    effects: list,
    ci_lower: list,
    ci_upper: list,
    weights: list = None,
    overall_effect: float = None,
    overall_ci: tuple = None,
    title: str = "Forest Plot — Meta-Analyse",
    xlabel: str = "Effektgröße (Standardisierte Mittlere Differenz)",
    output_path: str = "forest_plot.png",
    subtitle: str = None,
    source: str = None,
):
    """
    Erstellt einen publikationsreifen Forest Plot.

    Args:
        studies: Liste von Studiennamen
        effects: Liste von Effektgrößen
        ci_lower: Liste unterer Konfidenzintervall-Grenzen
        ci_upper: Liste oberer Konfidenzintervall-Grenzen
        weights: Optionale Gewichte (bestimmen Punktgröße)
        overall_effect: Gesamteffekt (Diamant)
        overall_ci: Konfidenzintervall des Gesamteffekts (lower, upper)
        title: Titel des Plots
        xlabel: Beschriftung der x-Achse
        output_path: Speicherpfad
        subtitle: Optionaler Untertitel
        source: Quellenangabe
    """
    n = len(studies)
    if weights is None:
        weights = [1.0] * n

    max_weight = max(weights) if weights else 1
    normalized_weights = [w / max_weight for w in weights]

    fig_height = max(6, n * 0.65 + 3)
    fig, ax = plt.subplots(figsize=(14, fig_height))

    y_positions = list(range(n - 1, -1, -1))

    # Null-Linie
    ax.axvline(x=0, color=COLORS['text_light'], linestyle='-', linewidth=1.2, alpha=0.7, zorder=1)

    # Studien plotten
    for i, (study, effect, cl, cu, w, y) in enumerate(
        zip(studies, effects, ci_lower, ci_upper, normalized_weights, y_positions)
    ):
        # Konfidenzintervall-Linie
        ax.plot([cl, cu], [y, y], color=COLORS['primary_dark'], linewidth=1.8, zorder=2)
        # Endkappen
        cap_height = 0.15
        ax.plot([cl, cl], [y - cap_height, y + cap_height], color=COLORS['primary_dark'], linewidth=1.5, zorder=2)
        ax.plot([cu, cu], [y - cap_height, y + cap_height], color=COLORS['primary_dark'], linewidth=1.5, zorder=2)
        # Punkt (Größe proportional zum Gewicht)
        marker_size = 80 + w * 220
        color = COLORS['significant'] if (cl > 0 or cu < 0) else COLORS['primary']
        ax.scatter(effect, y, s=marker_size, color=color, edgecolors='white',
                   linewidths=1.5, zorder=3, marker='s')

    # Gesamteffekt (Diamant)
    if overall_effect is not None and overall_ci is not None:
        diamond_y = -1.5
        diamond_x = [overall_ci[0], overall_effect, overall_ci[1], overall_effect]
        diamond_y_pts = [diamond_y, diamond_y + 0.4, diamond_y, diamond_y - 0.4]
        ax.fill(diamond_x, diamond_y_pts, color=COLORS['dopamin'], alpha=0.85, zorder=3)
        ax.plot(diamond_x + [diamond_x[0]], diamond_y_pts + [diamond_y_pts[0]],
                color=COLORS['dopamin'], linewidth=1.5, zorder=3)
        ax.text(-0.02, diamond_y, 'Gesamteffekt', ha='right', va='center',
                fontsize=12, fontweight='bold', color=COLORS['text_dark'],
                transform=ax.get_yaxis_transform())

    # Y-Achse
    all_y = y_positions + ([-1.5] if overall_effect is not None else [])
    ax.set_yticks(y_positions)
    ax.set_yticklabels(studies, fontsize=11)
    ax.set_ylim(min(all_y) - 1, max(y_positions) + 1)

    # Achsenbeschriftung
    ax.set_xlabel(xlabel, fontsize=13, color=COLORS['text_dark'], labelpad=10)

    # Titel
    title_y = 1.02
    ax.set_title(title, fontsize=20, fontweight='bold', color=COLORS['primary_dark'],
                 pad=20, loc='left')
    if subtitle:
        ax.text(0, 1.01, subtitle, transform=ax.transAxes, fontsize=11,
                color=COLORS['text_light'], style='italic')

    # Werte rechts anzeigen
    x_max = max(ci_upper) * 1.15
    for i, (effect, cl, cu, w_raw, y) in enumerate(
        zip(effects, ci_lower, ci_upper, weights, y_positions)
    ):
        text = f"{effect:.2f} [{cl:.2f}, {cu:.2f}]"
        ax.text(x_max * 0.95, y, text, ha='right', va='center', fontsize=9,
                fontfamily='monospace', color=COLORS['text_dark'],
                bbox=dict(boxstyle='round,pad=0.3', facecolor=COLORS['background_alt'],
                          edgecolor='none', alpha=0.8))

    # Favours-Labels
    xlim = ax.get_xlim()
    ax.text(xlim[0] * 0.5, min(all_y) - 0.7, '← Favours Control',
            ha='center', fontsize=9, color=COLORS['text_light'], style='italic')
    ax.text(xlim[1] * 0.5, min(all_y) - 0.7, 'Favours Treatment →',
            ha='center', fontsize=9, color=COLORS['text_light'], style='italic')

    # Quelle
    if source:
        fig.text(0.99, 0.01, f"Quelle: {source}", ha='right', va='bottom',
                 fontsize=8, color=COLORS['text_light'], style='italic')

    ax.spines['left'].set_visible(False)
    ax.tick_params(axis='y', length=0)

    plt.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Forest Plot gespeichert: {output_path}")
    return output_path


# =============================================================================
# Neurotransmitter-Pathway-Diagramm
# =============================================================================

def create_neurotransmitter_pathway(
    transmitter: str,
    synthesis_steps: list = None,
    receptors: list = None,
    brain_regions: list = None,
    effects: list = None,
    title: str = None,
    output_path: str = "pathway.png",
    source: str = None,
):
    """
    Erstellt ein professionelles Neurotransmitter-Pathway-Diagramm.

    Args:
        transmitter: Name des Neurotransmitters
        synthesis_steps: Schritte der Synthese
        receptors: Beteiligte Rezeptoren
        brain_regions: Involvierte Hirnregionen
        effects: Klinische/physiologische Effekte
        title: Optionaler Titel
        output_path: Speicherpfad
        source: Quellenangabe
    """
    if title is None:
        title = f"{transmitter}-System — Signalweg & Wirkung"

    color = NEUROTRANSMITTER_COLORS.get(transmitter, COLORS['primary'])

    fig, ax = plt.subplots(figsize=(18, 12), facecolor='white')
    ax.set_xlim(0, 18)
    ax.set_ylim(0, 12)
    ax.axis('off')

    # Titel
    ax.text(9, 11.5, title, ha='center', va='center', fontsize=24, fontweight='bold',
            color=COLORS['primary_dark'],
            path_effects=[pe.withStroke(linewidth=0, foreground='white')])

    def draw_box(x, y, w, h, text, box_color, text_color='white', fontsize=11,
                 alpha=0.9, rounded=True, shadow=True):
        """Zeichnet eine professionelle Box mit optionalem Schatten."""
        if shadow:
            shadow_box = FancyBboxPatch(
                (x + 0.05, y - 0.05), w, h,
                boxstyle="round,pad=0.15" if rounded else "square,pad=0.05",
                facecolor='#00000015', edgecolor='none', zorder=1
            )
            ax.add_patch(shadow_box)

        box = FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.15" if rounded else "square,pad=0.05",
            facecolor=box_color, edgecolor='white', linewidth=2,
            alpha=alpha, zorder=2
        )
        ax.add_patch(box)
        ax.text(x + w / 2, y + h / 2, text, ha='center', va='center',
                fontsize=fontsize, color=text_color, fontweight='bold',
                zorder=3, wrap=True)

    def draw_arrow(x1, y1, x2, y2, arrow_color=COLORS['text_light'], style='->', lw=2):
        """Zeichnet einen eleganten Pfeil."""
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                     arrowprops=dict(arrowstyle=style, color=arrow_color,
                                     lw=lw, connectionstyle='arc3,rad=0.0'),
                     zorder=2)

    # Zentraler Transmitter-Kreis
    circle = Circle((9, 7), 1.2, facecolor=color, edgecolor='white',
                     linewidth=3, alpha=0.95, zorder=5)
    ax.add_patch(circle)
    ax.text(9, 7, transmitter, ha='center', va='center', fontsize=18,
            fontweight='bold', color='white', zorder=6)

    # Synthese (links)
    if synthesis_steps:
        ax.text(2.5, 10.5, "Synthese", ha='center', fontsize=14, fontweight='bold',
                color=COLORS['text_dark'])
        for i, step in enumerate(synthesis_steps[:5]):
            y_pos = 9.5 - i * 1.2
            draw_box(0.5, y_pos - 0.3, 4, 0.6, step,
                     box_color=color, alpha=0.6 + i * 0.08, fontsize=10)
            if i < len(synthesis_steps) - 1:
                draw_arrow(2.5, y_pos - 0.3, 2.5, y_pos - 0.9, arrow_color=color)
        # Pfeil von Synthese zum Zentrum
        draw_arrow(4.5, 7, 7.8, 7, arrow_color=color, lw=2.5)

    # Rezeptoren (rechts)
    if receptors:
        ax.text(15.5, 10.5, "Rezeptoren", ha='center', fontsize=14, fontweight='bold',
                color=COLORS['text_dark'])
        for i, receptor in enumerate(receptors[:6]):
            y_pos = 9.5 - i * 1.0
            draw_box(13.5, y_pos - 0.25, 4, 0.5, receptor,
                     box_color=COLORS['primary_dark'], alpha=0.7 + i * 0.05,
                     fontsize=10)
        # Pfeil vom Zentrum zu Rezeptoren
        draw_arrow(10.2, 7, 13.5, 7, arrow_color=COLORS['primary_dark'], lw=2.5)

    # Hirnregionen (oben)
    if brain_regions:
        ax.text(9, 10.8, "Hirnregionen", ha='center', fontsize=14, fontweight='bold',
                color=COLORS['text_dark'])
        n_regions = min(len(brain_regions), 5)
        total_width = n_regions * 3
        start_x = 9 - total_width / 2
        for i, region in enumerate(brain_regions[:5]):
            x_pos = start_x + i * 3
            region_color = COLORS.get(region.lower().replace(' ', '_').replace('-', '_'),
                                      COLORS['primary_light'])
            draw_box(x_pos, 9.8, 2.6, 0.6, region, box_color=region_color, fontsize=9)
            draw_arrow(x_pos + 1.3, 9.8, 9, 8.2, arrow_color=region_color, lw=1.5)

    # Effekte (unten)
    if effects:
        ax.text(9, 3.8, "Klinische Effekte", ha='center', fontsize=14, fontweight='bold',
                color=COLORS['text_dark'])
        n_effects = min(len(effects), 5)
        total_width = n_effects * 3.2
        start_x = 9 - total_width / 2
        for i, effect in enumerate(effects[:5]):
            x_pos = start_x + i * 3.2
            draw_box(x_pos, 2.5, 2.8, 0.8, effect,
                     box_color=COLORS['accent'], text_color=COLORS['text_dark'],
                     fontsize=9, alpha=0.85)
            draw_arrow(9, 5.8, x_pos + 1.4, 3.3, arrow_color=COLORS['accent'], lw=1.5)

    # Quelle
    if source:
        ax.text(17.5, 0.3, f"Quelle: {source}", ha='right', va='bottom',
                fontsize=8, color=COLORS['text_light'], style='italic')

    # NEVPAZ Branding
    ax.text(0.5, 0.3, "NEVPAZ Scientific", ha='left', va='bottom',
            fontsize=9, color=COLORS['primary'], fontweight='bold', alpha=0.6)

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Pathway-Diagramm gespeichert: {output_path}")
    return output_path


# =============================================================================
# Rezeptor-Bindungsprofil (Spider/Radar Chart)
# =============================================================================

def create_receptor_binding_profile(
    drug_name: str,
    receptors: dict,
    comparison_drug: str = None,
    comparison_receptors: dict = None,
    title: str = None,
    output_path: str = "binding_profile.png",
    source: str = None,
):
    """
    Erstellt ein Rezeptor-Bindungsprofil als Radar-Chart.

    Args:
        drug_name: Name des Medikaments
        receptors: Dict {Rezeptorname: Ki-Wert in nM}
        comparison_drug: Optionaler Vergleichsmedikament-Name
        comparison_receptors: Optionale Ki-Werte des Vergleichsmedikaments
        title: Optionaler Titel
        output_path: Speicherpfad
        source: Quellenangabe
    """
    if title is None:
        title = f"Rezeptor-Bindungsprofil — {drug_name}"

    labels = list(receptors.keys())
    n = len(labels)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
    angles += angles[:1]

    # Ki-Werte logarithmisch transformieren (niedrig = hohe Affinität)
    max_ki = max(receptors.values()) * 1.5
    values = [1 - (np.log10(v + 1) / np.log10(max_ki + 1)) for v in receptors.values()]
    values += values[:1]

    fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(polar=True), facecolor='white')

    # Hintergrund-Ringe
    ax.set_facecolor(COLORS['background_alt'])
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)

    # Hauptdaten
    ax.fill(angles, values, color=COLORS['dopamin'], alpha=0.2)
    ax.plot(angles, values, color=COLORS['dopamin'], linewidth=2.5,
            label=drug_name, marker='o', markersize=8)

    # Vergleichsmedikament
    if comparison_drug and comparison_receptors:
        comp_values = [1 - (np.log10(comparison_receptors.get(l, max_ki) + 1) / np.log10(max_ki + 1))
                       for l in labels]
        comp_values += comp_values[:1]
        ax.fill(angles, comp_values, color=COLORS['serotonin'], alpha=0.15)
        ax.plot(angles, comp_values, color=COLORS['serotonin'], linewidth=2,
                label=comparison_drug, marker='s', markersize=7, linestyle='--')

    # Achsenbeschriftung
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=11, fontweight='bold', color=COLORS['text_dark'])

    # Y-Achse
    ax.set_ylim(0, 1)
    ax.set_yticks([0.25, 0.5, 0.75, 1.0])
    ax.set_yticklabels(['Niedrig', 'Mittel', 'Hoch', 'Sehr hoch'],
                       fontsize=8, color=COLORS['text_light'])

    # Titel & Legende
    ax.set_title(title, fontsize=18, fontweight='bold', color=COLORS['primary_dark'],
                 pad=30)
    ax.legend(loc='lower right', bbox_to_anchor=(1.15, -0.05), fontsize=11,
              framealpha=0.9, edgecolor=COLORS['grid'])

    # Quelle
    if source:
        fig.text(0.95, 0.02, f"Quelle: {source}", ha='right', va='bottom',
                 fontsize=8, color=COLORS['text_light'], style='italic')

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Bindungsprofil gespeichert: {output_path}")
    return output_path


# =============================================================================
# Dosis-Wirkungs-Kurve
# =============================================================================

def create_dose_response_curve(
    drugs: list,
    ec50_values: list,
    emax_values: list = None,
    hill_coefficients: list = None,
    dose_range: tuple = (0.001, 1000),
    title: str = "Dosis-Wirkungs-Kurve",
    xlabel: str = "Konzentration (nM)",
    ylabel: str = "Wirkung (%)",
    output_path: str = "dose_response.png",
    source: str = None,
):
    """
    Erstellt Dosis-Wirkungs-Kurven nach dem Hill-Modell.

    Args:
        drugs: Liste von Medikamentennamen
        ec50_values: EC50-Werte
        emax_values: Maximale Wirkung (Default: 100%)
        hill_coefficients: Hill-Koeffizienten (Default: 1)
        dose_range: Dosisbereich (min, max)
        title: Titel
        xlabel: X-Achsenbeschriftung
        ylabel: Y-Achsenbeschriftung
        output_path: Speicherpfad
        source: Quellenangabe
    """
    n = len(drugs)
    if emax_values is None:
        emax_values = [100.0] * n
    if hill_coefficients is None:
        hill_coefficients = [1.0] * n

    doses = np.logspace(np.log10(dose_range[0]), np.log10(dose_range[1]), 500)
    drug_colors = [COLORS['dopamin'], COLORS['serotonin'], COLORS['noradrenalin'],
                   COLORS['gaba'], COLORS['glutamat'], COLORS['acetylcholin']]

    fig, ax = plt.subplots(figsize=(12, 8), facecolor='white')

    for i, (drug, ec50, emax, hill) in enumerate(
        zip(drugs, ec50_values, emax_values, hill_coefficients)
    ):
        response = emax * (doses ** hill) / (ec50 ** hill + doses ** hill)
        c = drug_colors[i % len(drug_colors)]
        ax.semilogx(doses, response, linewidth=2.5, color=c, label=drug, zorder=3)

        # EC50-Markierung
        ec50_response = emax * 0.5
        ax.plot([ec50, ec50], [0, ec50_response], '--', color=c, alpha=0.4, linewidth=1)
        ax.plot([dose_range[0], ec50], [ec50_response, ec50_response], '--', color=c, alpha=0.4, linewidth=1)
        ax.scatter([ec50], [ec50_response], color=c, s=60, zorder=4, edgecolors='white', linewidths=1.5)
        ax.annotate(f'EC₅₀={ec50:.1f}', xy=(ec50, ec50_response),
                    xytext=(ec50 * 2, ec50_response + 5),
                    fontsize=9, color=c, fontweight='bold',
                    arrowprops=dict(arrowstyle='->', color=c, lw=1))

    ax.set_xlabel(xlabel, fontsize=13, color=COLORS['text_dark'])
    ax.set_ylabel(ylabel, fontsize=13, color=COLORS['text_dark'])
    ax.set_title(title, fontsize=20, fontweight='bold', color=COLORS['primary_dark'], pad=15)
    ax.set_ylim(-5, max(emax_values) * 1.1)
    ax.legend(fontsize=11, loc='lower right', framealpha=0.9, edgecolor=COLORS['grid'])

    if source:
        fig.text(0.99, 0.01, f"Quelle: {source}", ha='right', va='bottom',
                 fontsize=8, color=COLORS['text_light'], style='italic')

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Dosis-Wirkungs-Kurve gespeichert: {output_path}")
    return output_path


# =============================================================================
# Synaptischer Spalt — Detaildiagramm
# =============================================================================

def create_synaptic_cleft_diagram(
    transmitter: str = "Dopamin",
    title: str = None,
    output_path: str = "synaptic_cleft.png",
    show_reuptake: bool = True,
    show_degradation: bool = True,
    show_autoreceptors: bool = True,
    drug_mechanism: str = None,
    source: str = None,
):
    """
    Erstellt ein detailliertes Diagramm des synaptischen Spalts.

    Args:
        transmitter: Neurotransmitter-Name
        title: Optionaler Titel
        output_path: Speicherpfad
        show_reuptake: Wiederaufnahme-Transporter anzeigen
        show_degradation: Abbaumechanismen anzeigen
        show_autoreceptors: Autorezeptoren anzeigen
        drug_mechanism: Optionaler Wirkmechanismus eines Medikaments
        source: Quellenangabe
    """
    if title is None:
        title = f"Synaptischer Spalt — {transmitter}-System"

    color = NEUROTRANSMITTER_COLORS.get(transmitter, COLORS['primary'])

    fig, ax = plt.subplots(figsize=(16, 12), facecolor='white')
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 12)
    ax.axis('off')

    # Titel
    ax.text(8, 11.5, title, ha='center', va='center', fontsize=22,
            fontweight='bold', color=COLORS['primary_dark'])

    # Präsynaptische Membran
    membrane_pre = FancyBboxPatch((1, 7.5), 14, 2.5,
                                   boxstyle="round,pad=0.3",
                                   facecolor='#FDE8D0', edgecolor='#E67E22',
                                   linewidth=3, alpha=0.8)
    ax.add_patch(membrane_pre)
    ax.text(8, 9.5, "Präsynaptisches Neuron", ha='center', fontsize=14,
            fontweight='bold', color=COLORS['text_dark'])

    # Postsynaptische Membran
    membrane_post = FancyBboxPatch((1, 1), 14, 2.5,
                                    boxstyle="round,pad=0.3",
                                    facecolor='#D5E8D4', edgecolor='#2ECC71',
                                    linewidth=3, alpha=0.8)
    ax.add_patch(membrane_post)
    ax.text(8, 2.8, "Postsynaptisches Neuron", ha='center', fontsize=14,
            fontweight='bold', color=COLORS['text_dark'])

    # Synaptischer Spalt Label
    ax.text(8, 5.5, "Synaptischer Spalt", ha='center', fontsize=12,
            color=COLORS['text_light'], style='italic',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='white',
                      edgecolor=COLORS['grid'], alpha=0.9))

    # Vesikel
    for vx, vy in [(4, 8.8), (6, 9.0), (8, 8.7), (10, 8.9), (12, 8.8)]:
        vesicle = Circle((vx, vy), 0.35, facecolor=color, edgecolor='white',
                         linewidth=2, alpha=0.7, zorder=5)
        ax.add_patch(vesicle)

    # Transmitter-Moleküle im Spalt
    np.random.seed(42)
    for _ in range(15):
        mx = np.random.uniform(2, 14)
        my = np.random.uniform(4, 7)
        molecule = Circle((mx, my), 0.12, facecolor=color, edgecolor='white',
                          linewidth=1, alpha=0.8, zorder=4)
        ax.add_patch(molecule)
    ax.text(14, 6.8, f"{transmitter}\nMoleküle", ha='center', fontsize=9,
            color=color, fontweight='bold')

    # Postsynaptische Rezeptoren
    receptor_positions = [3, 5, 7, 9, 11, 13]
    receptor_types = ['D1', 'D2', 'D3', 'D4', 'D1', 'D2'] if transmitter == 'Dopamin' else \
                     [f'R{i+1}' for i in range(6)]
    for rx, rtype in zip(receptor_positions, receptor_types):
        rect = FancyBboxPatch((rx - 0.4, 3.2), 0.8, 0.8,
                               boxstyle="round,pad=0.1",
                               facecolor=COLORS['primary_dark'], edgecolor='white',
                               linewidth=2, alpha=0.85, zorder=5)
        ax.add_patch(rect)
        ax.text(rx, 3.6, rtype, ha='center', va='center', fontsize=9,
                fontweight='bold', color='white', zorder=6)

    # Wiederaufnahme-Transporter
    if show_reuptake:
        for tx in [3.5, 8, 12.5]:
            trans = FancyBboxPatch((tx - 0.5, 7.5), 1, 0.8,
                                    boxstyle="round,pad=0.1",
                                    facecolor='#E74C3C', edgecolor='white',
                                    linewidth=2, alpha=0.8, zorder=5)
            ax.add_patch(trans)
            ax.text(tx, 7.9, 'DAT' if transmitter == 'Dopamin' else 'T',
                    ha='center', va='center', fontsize=8,
                    fontweight='bold', color='white', zorder=6)
        ax.text(14.5, 7.5, "Transporter\n(Reuptake)", ha='center', fontsize=9,
                color='#E74C3C', fontweight='bold')

    # Autorezeptoren
    if show_autoreceptors:
        for arx in [5, 10]:
            auto = FancyBboxPatch((arx - 0.35, 7.6), 0.7, 0.7,
                                   boxstyle="round,pad=0.1",
                                   facecolor='#F39C12', edgecolor='white',
                                   linewidth=1.5, alpha=0.8, zorder=5)
            ax.add_patch(auto)
            ax.text(arx, 7.95, 'Auto', ha='center', va='center', fontsize=7,
                    fontweight='bold', color='white', zorder=6)

    # Degradation (MAO/COMT)
    if show_degradation:
        ax.text(2, 5, "MAO", ha='center', fontsize=10, fontweight='bold',
                color='#8E44AD',
                bbox=dict(boxstyle='round,pad=0.3', facecolor='#D7BDE2',
                          edgecolor='#8E44AD', alpha=0.8))
        ax.text(2, 4.2, "COMT", ha='center', fontsize=10, fontweight='bold',
                color='#8E44AD',
                bbox=dict(boxstyle='round,pad=0.3', facecolor='#D7BDE2',
                          edgecolor='#8E44AD', alpha=0.8))

    # Medikamenten-Mechanismus
    if drug_mechanism:
        ax.text(8, 0.4, f"Wirkmechanismus: {drug_mechanism}", ha='center',
                fontsize=11, fontweight='bold', color=COLORS['dopamin'],
                bbox=dict(boxstyle='round,pad=0.5', facecolor='#FADBD8',
                          edgecolor=COLORS['dopamin'], alpha=0.9))

    # Legende
    legend_items = [
        (color, f"{transmitter}-Vesikel/Moleküle"),
        (COLORS['primary_dark'], "Postsynaptische Rezeptoren"),
    ]
    if show_reuptake:
        legend_items.append(('#E74C3C', "Wiederaufnahme-Transporter"))
    if show_autoreceptors:
        legend_items.append(('#F39C12', "Autorezeptoren"))
    if show_degradation:
        legend_items.append(('#8E44AD', "Enzyme (MAO/COMT)"))

    for i, (lc, lt) in enumerate(legend_items):
        y_leg = 1.8 - i * 0.35
        circle_leg = Circle((1.5, y_leg), 0.12, facecolor=lc, edgecolor='white',
                             linewidth=1, zorder=5)
        ax.add_patch(circle_leg)
        ax.text(1.8, y_leg, lt, va='center', fontsize=9, color=COLORS['text_dark'])

    if source:
        ax.text(15.5, 0.3, f"Quelle: {source}", ha='right', va='bottom',
                fontsize=8, color=COLORS['text_light'], style='italic')

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Synaptischer-Spalt-Diagramm gespeichert: {output_path}")
    return output_path


# =============================================================================
# PPTX-Erstellung
# =============================================================================

def create_scientific_pptx(
    title: str,
    subtitle: str,
    slides_data: list,
    output_path: str = "presentation.pptx",
    author: str = "NEVPAZ Praxis",
    image_paths: list = None,
):
    """
    Erstellt eine professionelle wissenschaftliche PowerPoint-Präsentation.

    Args:
        title: Präsentationstitel
        subtitle: Untertitel
        slides_data: Liste von Dicts mit Foliendaten:
            [{"title": "...", "content": "...", "image": "path.png", "layout": "title_content"}]
        output_path: Speicherpfad
        author: Autorenname
        image_paths: Optionale Liste von Bildpfaden
    """
    if not HAS_PPTX:
        print("FEHLER: python-pptx nicht installiert. pip install python-pptx")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Farben
    primary_dark = RGBColor(0x1a, 0x52, 0x76)
    primary = RGBColor(0x29, 0x80, 0xb9)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    text_dark = RGBColor(0x2C, 0x3E, 0x50)
    bg_alt = RGBColor(0xF8, 0xF9, 0xFA)

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Hintergrund
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = primary_dark

        # Titel
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = white
        p.alignment = PP_ALIGN.LEFT

        # Untertitel
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x5D, 0xAD, 0xE2)
        p2.alignment = PP_ALIGN.LEFT

        # Linie
        from pptx.shapes.autoshape import Shape
        line = slide.shapes.add_connector(1, Inches(1), Inches(4), Inches(5), Inches(4))
        line.line.color.rgb = RGBColor(0xF3, 0x9C, 0x12)
        line.line.width = Pt(3)

        # Autor/Datum
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = author
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        p3.alignment = PP_ALIGN.LEFT

    def add_content_slide(slide_title, content_text, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Header-Leiste
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = primary_dark
        header.line.fill.background()

        # Titel in Header
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = white
        p.alignment = PP_ALIGN.LEFT

        if image_path and os.path.exists(image_path):
            # Layout mit Bild
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5),
                                      width=Inches(7.5))
            # Text rechts
            txBox2 = slide.shapes.add_textbox(Inches(8.3), Inches(1.5),
                                               Inches(4.5), Inches(5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for line in content_text.split('\n'):
                p2 = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
                p2.text = line
                p2.font.size = Pt(14)
                p2.font.color.rgb = text_dark
                p2.space_after = Pt(8)
        else:
            # Nur Text
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
                                               Inches(11.7), Inches(5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for line in content_text.split('\n'):
                p2 = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
                p2.text = line
                p2.font.size = Pt(16)
                p2.font.color.rgb = text_dark
                p2.space_after = Pt(10)

        # Footer
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8),
                                           Inches(12), Inches(0.4))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"{author} | {slide_title}"
        p3.font.size = Pt(9)
        p3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        p3.alignment = PP_ALIGN.RIGHT

    # Titelfolie
    add_title_slide(title, subtitle)

    # Inhaltsfolien
    for slide_data in slides_data:
        add_content_slide(
            slide_data.get('title', 'Folie'),
            slide_data.get('content', ''),
            slide_data.get('image', None)
        )

    prs.save(output_path)
    print(f"Präsentation gespeichert: {output_path}")
    return output_path


# =============================================================================
# Interaktive Plotly-Grafiken
# =============================================================================

def create_interactive_heatmap(
    data_matrix: list,
    x_labels: list,
    y_labels: list,
    title: str = "Korrelationsmatrix",
    output_path: str = "heatmap.html",
    colorscale: str = "RdBu_r",
    source: str = None,
):
    """
    Erstellt eine interaktive Heatmap mit Plotly.

    Args:
        data_matrix: 2D-Liste oder numpy-Array mit Werten
        x_labels: Spaltenbezeichnungen
        y_labels: Zeilenbezeichnungen
        title: Titel
        output_path: Speicherpfad (.html für interaktiv, .png für statisch)
        colorscale: Plotly-Farbskala
        source: Quellenangabe
    """
    if not HAS_PLOTLY:
        print("FEHLER: plotly nicht installiert. pip install plotly kaleido")
        return None

    fig = go.Figure(data=go.Heatmap(
        z=data_matrix,
        x=x_labels,
        y=y_labels,
        colorscale=colorscale,
        text=[[f"{v:.2f}" for v in row] for row in data_matrix],
        texttemplate="%{text}",
        textfont={"size": 12},
        hovertemplate="<b>%{x}</b> vs <b>%{y}</b><br>Wert: %{z:.3f}<extra></extra>",
    ))

    annotation_text = f"<br><i style='font-size:10px'>Quelle: {source}</i>" if source else ""

    fig.update_layout(
        title=dict(
            text=f"<b>{title}</b>{annotation_text}",
            font=dict(size=20, color=COLORS['primary_dark']),
            x=0.5,
        ),
        width=800,
        height=700,
        font=dict(family="Calibri, Arial, sans-serif"),
        paper_bgcolor='white',
        plot_bgcolor='white',
    )

    if output_path.endswith('.html'):
        fig.write_html(output_path)
        print(f"Interaktive Heatmap gespeichert: {output_path}")
        # Auch als PNG
        png_path = output_path.replace('.html', '.png')
        try:
            fig.write_image(png_path, scale=3)
            print(f"Statische Heatmap gespeichert: {png_path}")
        except Exception:
            pass
    else:
        try:
            fig.write_image(output_path, scale=3)
            print(f"Heatmap gespeichert: {output_path}")
        except Exception as e:
            print(f"Fehler beim Speichern: {e}")

    return output_path


# =============================================================================
# Studienvergleichs-Tabelle
# =============================================================================

def create_study_comparison_table(
    studies: list,
    columns: list = None,
    title: str = "Studienvergleich",
    output_path: str = "study_comparison.png",
    source: str = None,
):
    """
    Erstellt eine professionelle Vergleichstabelle mehrerer Studien.

    Args:
        studies: Liste von Dicts mit Studiendaten
            [{"Studie": "...", "N": 100, "Design": "RCT", ...}]
        columns: Spalten die angezeigt werden sollen (Default: alle Keys)
        title: Tabellentitel
        output_path: Speicherpfad
        source: Quellenangabe
    """
    if columns is None:
        columns = list(studies[0].keys())

    n_rows = len(studies)
    n_cols = len(columns)

    fig_width = max(12, n_cols * 2.2)
    fig_height = max(4, n_rows * 0.6 + 2.5)

    fig, ax = plt.subplots(figsize=(fig_width, fig_height), facecolor='white')
    ax.axis('off')

    # Titel
    ax.text(0.5, 0.98, title, transform=ax.transAxes, ha='center', va='top',
            fontsize=20, fontweight='bold', color=COLORS['primary_dark'])

    # Tabellendaten
    cell_text = [[str(study.get(col, '—')) for col in columns] for study in studies]

    table = ax.table(
        cellText=cell_text,
        colLabels=columns,
        cellLoc='center',
        loc='center',
        bbox=[0.02, 0.05, 0.96, 0.85]
    )

    table.auto_set_font_size(False)
    table.set_fontsize(10)

    # Header-Stil
    for j in range(n_cols):
        cell = table[0, j]
        cell.set_facecolor(COLORS['primary_dark'])
        cell.set_text_props(color='white', fontweight='bold', fontsize=11)
        cell.set_edgecolor('white')
        cell.set_linewidth(2)

    # Zeilen-Stil (alternierend)
    for i in range(1, n_rows + 1):
        for j in range(n_cols):
            cell = table[i, j]
            cell.set_facecolor(COLORS['background'] if i % 2 == 0 else COLORS['background_alt'])
            cell.set_edgecolor(COLORS['grid'])
            cell.set_linewidth(0.5)
            cell.set_text_props(fontsize=10)

    # Quelle
    if source:
        fig.text(0.98, 0.01, f"Quelle: {source}", ha='right', va='bottom',
                 fontsize=8, color=COLORS['text_light'], style='italic')

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"Studienvergleich gespeichert: {output_path}")
    return output_path


# =============================================================================
# PRISMA Flussdiagramm
# =============================================================================

def create_prisma_flowchart(
    identified: int,
    screened: int,
    eligible: int,
    included: int,
    excluded_reasons: dict = None,
    title: str = "PRISMA Flussdiagramm",
    output_path: str = "prisma.png",
    source: str = None,
):
    """
    Erstellt ein PRISMA-Flussdiagramm für systematische Reviews.

    Args:
        identified: Anzahl identifizierter Artikel
        screened: Anzahl gescreenter Artikel
        eligible: Anzahl geprüfter Artikel
        included: Anzahl eingeschlossener Artikel
        excluded_reasons: Dict mit Ausschlussgründen {Grund: Anzahl}
        title: Titel
        output_path: Speicherpfad
        source: Quellenangabe
    """
    fig, ax = plt.subplots(figsize=(14, 16), facecolor='white')
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 16)
    ax.axis('off')

    # Titel
    ax.text(7, 15.5, title, ha='center', va='center', fontsize=22,
            fontweight='bold', color=COLORS['primary_dark'])

    def prisma_box(x, y, w, h, text, color=COLORS['primary'], text_color='white'):
        box = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.2",
                              facecolor=color, edgecolor='white', linewidth=2, alpha=0.9)
        ax.add_patch(box)
        ax.text(x + w/2, y + h/2, text, ha='center', va='center',
                fontsize=12, color=text_color, fontweight='bold', wrap=True)

    def arrow_down(x, y1, y2):
        ax.annotate('', xy=(x, y2), xytext=(x, y1),
                     arrowprops=dict(arrowstyle='->', color=COLORS['text_dark'],
                                     lw=2.5))

    def arrow_right(x1, y, x2):
        ax.annotate('', xy=(x2, y), xytext=(x1, y),
                     arrowprops=dict(arrowstyle='->', color=COLORS['not_significant'],
                                     lw=2))

    # Identifikation
    prisma_box(3.5, 13.5, 7, 1.2,
               f"Identifizierte Artikel\n(n = {identified})",
               color=COLORS['primary_dark'])

    arrow_down(7, 13.5, 12.2)

    # Duplikate entfernt
    duplicates = identified - screened
    prisma_box(3.5, 11, 7, 1.2,
               f"Nach Duplikatentfernung\n(n = {screened})",
               color=COLORS['primary'])
    arrow_right(10.5, 11.6, 11)
    prisma_box(11, 11, 2.5, 1.2,
               f"Duplikate\n(n = {duplicates})",
               color=COLORS['not_significant'])

    arrow_down(7, 11, 9.7)

    # Screening
    excluded_screening = screened - eligible
    prisma_box(3.5, 8.5, 7, 1.2,
               f"Gescreente Artikel\n(n = {screened})",
               color=COLORS['primary'])
    arrow_right(10.5, 9.1, 11)
    prisma_box(11, 8.5, 2.5, 1.2,
               f"Ausgeschlossen\n(n = {excluded_screening})",
               color=COLORS['not_significant'])

    arrow_down(7, 8.5, 7.2)

    # Eignung
    excluded_eligibility = eligible - included
    prisma_box(3.5, 6, 7, 1.2,
               f"Geprüfte Volltexte\n(n = {eligible})",
               color=COLORS['primary'])

    if excluded_reasons:
        reasons_text = f"Ausgeschlossen (n = {excluded_eligibility}):\n"
        reasons_text += "\n".join([f"• {k}: {v}" for k, v in excluded_reasons.items()])
        box_height = max(1.2, len(excluded_reasons) * 0.4 + 0.8)
        arrow_right(10.5, 6.6, 11)
        prisma_box(11, 6.6 - box_height/2, 2.8, box_height,
                   reasons_text, color=COLORS['not_significant'], text_color='white')
    else:
        arrow_right(10.5, 6.6, 11)
        prisma_box(11, 6, 2.5, 1.2,
                   f"Ausgeschlossen\n(n = {excluded_eligibility})",
                   color=COLORS['not_significant'])

    arrow_down(7, 6, 4.7)

    # Eingeschlossen
    prisma_box(3.5, 3.5, 7, 1.2,
               f"Eingeschlossene Studien\n(n = {included})",
               color=COLORS['significant'])

    # Phasen-Labels (links)
    phases = [
        (14.2, "Identifikation"),
        (11.6, "Screening"),
        (9.1, "Screening"),
        (6.6, "Eignung"),
        (4.1, "Einschluss"),
    ]
    for py, plabel in phases:
        ax.text(0.5, py, plabel, ha='center', va='center', fontsize=10,
                color=COLORS['text_light'], rotation=90, fontweight='bold')

    if source:
        ax.text(13.5, 0.5, f"Quelle: {source}", ha='right', va='bottom',
                fontsize=8, color=COLORS['text_light'], style='italic')

    plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"PRISMA-Flussdiagramm gespeichert: {output_path}")
    return output_path


# =============================================================================
# Demo
# =============================================================================

def run_demo(output_dir: str = "demo_output"):
    """Erstellt Demo-Grafiken aller verfügbaren Typen."""

    os.makedirs(output_dir, exist_ok=True)
    print("=" * 60)
    print("NEVPAZ Scientific Visualization Toolkit — Demo")
    print("=" * 60)

    # 1. Forest Plot
    print("\n[1/6] Forest Plot...")
    create_forest_plot(
        studies=["Schmidt et al. 2023", "Müller et al. 2022", "Wagner et al. 2024",
                 "Fischer et al. 2023", "Weber et al. 2021", "Bauer et al. 2024",
                 "Schneider et al. 2022"],
        effects=[0.45, 0.32, 0.61, 0.28, 0.55, 0.39, 0.50],
        ci_lower=[0.12, -0.05, 0.30, -0.08, 0.22, 0.10, 0.18],
        ci_upper=[0.78, 0.69, 0.92, 0.64, 0.88, 0.68, 0.82],
        weights=[15, 12, 18, 10, 16, 14, 13],
        overall_effect=0.44,
        overall_ci=(0.30, 0.58),
        title="Forest Plot — Methylphenidat vs. Placebo bei ADHS",
        subtitle="Primärer Endpunkt: ADHS-RS Gesamtscore",
        source="Meta-Analyse, NEVPAZ 2024",
        output_path=os.path.join(output_dir, "forest_plot.png")
    )

    # 2. Pathway-Diagramm
    print("[2/6] Dopamin-Pathway...")
    create_neurotransmitter_pathway(
        transmitter="Dopamin",
        synthesis_steps=["Tyrosin", "L-DOPA\n(Tyrosinhydroxylase)", "Dopamin\n(AADC)"],
        receptors=["D1 (exzitatorisch)", "D2 (inhibitorisch)", "D3", "D4", "D5"],
        brain_regions=["VTA", "Substantia Nigra", "Präfrontaler Kortex",
                       "Nucleus Accumbens", "Striatum"],
        effects=["Motivation", "Belohnung", "Motorik", "Kognition", "Aufmerksamkeit"],
        source="Stahl's Essential Neuropharmacology, 2024",
        output_path=os.path.join(output_dir, "dopamin_pathway.png")
    )

    # 3. Rezeptor-Bindungsprofil
    print("[3/6] Rezeptor-Bindungsprofil...")
    create_receptor_binding_profile(
        drug_name="Aripiprazol",
        receptors={
            "D2": 0.34, "D3": 0.8, "5-HT1A": 1.7, "5-HT2A": 3.4,
            "5-HT2C": 15, "H1": 61, "α1": 57, "α2": 74,
            "M1": 6780, "DAT": 520, "SERT": 98,
        },
        comparison_drug="Risperidon",
        comparison_receptors={
            "D2": 3.77, "D3": 10.5, "5-HT1A": 210, "5-HT2A": 0.17,
            "5-HT2C": 26, "H1": 2.6, "α1": 0.7, "α2": 8,
            "M1": 4000, "DAT": 5000, "SERT": 2700,
        },
        source="PDSP Ki Database, 2024",
        output_path=os.path.join(output_dir, "binding_profile.png")
    )

    # 4. Dosis-Wirkungs-Kurve
    print("[4/6] Dosis-Wirkungs-Kurve...")
    create_dose_response_curve(
        drugs=["Methylphenidat", "Amphetamin", "Atomoxetin"],
        ec50_values=[30, 10, 100],
        emax_values=[95, 100, 80],
        hill_coefficients=[1.2, 1.5, 0.9],
        title="Dosis-Wirkungs-Kurven — ADHS-Medikation",
        source="Klinische Pharmakologie, NEVPAZ 2024",
        output_path=os.path.join(output_dir, "dose_response.png")
    )

    # 5. Synaptischer Spalt
    print("[5/6] Synaptischer Spalt...")
    create_synaptic_cleft_diagram(
        transmitter="Dopamin",
        drug_mechanism="Methylphenidat: DAT-Blockade → erhöhte DA-Konzentration im Spalt",
        source="Volkow et al., Am J Psychiatry, 2012",
        output_path=os.path.join(output_dir, "synaptic_cleft.png")
    )

    # 6. PRISMA
    print("[6/6] PRISMA-Flussdiagramm...")
    create_prisma_flowchart(
        identified=1247,
        screened=892,
        eligible=156,
        included=42,
        excluded_reasons={
            "Kein RCT": 48,
            "Falsches Outcome": 32,
            "Pädiatrisch": 21,
            "Keine Volltexte": 13,
        },
        title="PRISMA — Systematisches Review: ADHS-Pharmakotherapie",
        source="NEVPAZ Systematic Review, 2024",
        output_path=os.path.join(output_dir, "prisma.png")
    )

    print("\n" + "=" * 60)
    print(f"Alle Demo-Grafiken gespeichert in: {output_dir}/")
    print("=" * 60)


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="NEVPAZ Scientific Visualization Toolkit",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Beispiele:
  python scientific_viz.py --demo
  python scientific_viz.py --type forest --output my_forest.png
  python scientific_viz.py --type pathway --transmitter Serotonin
        """
    )
    parser.add_argument('--demo', action='store_true', help='Erstellt Demo-Grafiken')
    parser.add_argument('--type', choices=['forest', 'pathway', 'binding', 'dose_response',
                                            'synapse', 'prisma', 'heatmap', 'comparison'],
                        help='Diagrammtyp')
    parser.add_argument('--output', default=None, help='Ausgabepfad')
    parser.add_argument('--output-dir', default='output', help='Ausgabeverzeichnis')
    parser.add_argument('--transmitter', default='Dopamin', help='Neurotransmitter')
    parser.add_argument('--title', default=None, help='Diagrammtitel')

    args = parser.parse_args()

    if args.demo:
        run_demo(args.output_dir)
    elif args.type:
        os.makedirs(args.output_dir, exist_ok=True)
        output = args.output or os.path.join(args.output_dir, f"{args.type}.png")

        if args.type == 'forest':
            print("Verwende Demo-Daten für Forest Plot...")
            create_forest_plot(
                studies=["Studie A", "Studie B", "Studie C"],
                effects=[0.5, 0.3, 0.7],
                ci_lower=[0.1, -0.1, 0.3],
                ci_upper=[0.9, 0.7, 1.1],
                overall_effect=0.5,
                overall_ci=(0.3, 0.7),
                title=args.title or "Forest Plot",
                output_path=output,
            )
        elif args.type == 'pathway':
            create_neurotransmitter_pathway(
                transmitter=args.transmitter,
                synthesis_steps=["Vorstufe", "Zwischenstufe", args.transmitter],
                receptors=["R1", "R2", "R3"],
                effects=["Effekt 1", "Effekt 2", "Effekt 3"],
                title=args.title,
                output_path=output,
            )
        elif args.type == 'synapse':
            create_synaptic_cleft_diagram(
                transmitter=args.transmitter,
                title=args.title,
                output_path=output,
            )
        elif args.type == 'prisma':
            create_prisma_flowchart(
                identified=500, screened=350, eligible=80, included=25,
                title=args.title or "PRISMA Flussdiagramm",
                output_path=output,
            )
        print("Fertig!")
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
